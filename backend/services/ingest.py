"""
Universal document ingestion pipeline.

Supported inputs (best-effort; any failure is reported in the KbDocument row):
  - PDF (text + OCR fallback per page when no text is found)
  - DOCX (Word)
  - PPTX (PowerPoint)
  - TXT / MD / LOG
  - XLSX / XLS / CSV (one chunk per row; header-aware)
  - Images (PNG/JPG/JPEG/WEBP/TIFF) via Tesseract OCR (eng + hin if installed)
  - SQL dumps (.sql): splits into statements; one chunk per statement

Every extracted Chunk is:
  1) cleaned & chunked into ~800-character pieces with ~120-char overlap
  2) enriched with extracted entities (phone/IMEI/IP/date/email/cell/money)
  3) embedded with the multilingual MiniLM model
  4) persisted as (KbDocument, KbChunk) rows

Call `ingest_document(...)` from an API endpoint. Heavy work runs synchronously
on the request thread; for very large files consider wrapping the call in a
FastAPI BackgroundTask (see main.py).
"""
from __future__ import annotations

import hashlib
import io
import os
import re
import uuid
from datetime import datetime
from typing import Any, Dict, Iterable, List, Optional, Tuple

from sqlalchemy.orm import Session

from ..models import KbChunk, KbDocument
from . import embeddings as emb_svc
from .entity_extract import detect_language, extract_entities

CHUNK_SIZE = 800
CHUNK_OVERLAP = 120
MAX_CHUNKS_PER_FILE = int(os.getenv("DIP_KB_MAX_CHUNKS", "5000"))


# ---------------------------------------------------------------------------
# File-type detection
# ---------------------------------------------------------------------------


def _ext(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower().lstrip(".")


def _source_type(filename: str) -> str:
    e = _ext(filename)
    if e in {"pdf"}:
        return "pdf"
    if e in {"docx", "doc"}:
        return "docx"
    if e in {"pptx", "ppt"}:
        return "pptx"
    if e in {"xlsx", "xls"}:
        return "xlsx"
    if e in {"csv", "tsv"}:
        return "csv"
    if e in {"txt", "md", "log", "rtf"}:
        return "txt"
    if e in {"png", "jpg", "jpeg", "webp", "tiff", "bmp"}:
        return "image"
    if e in {"sql"}:
        return "sql"
    return "unknown"


# ---------------------------------------------------------------------------
# Parsers — each returns a list of (section, page_or_row, text_piece)
# ---------------------------------------------------------------------------


def _parse_txt(content: bytes) -> List[Tuple[Optional[str], Optional[int], str]]:
    try:
        text = content.decode("utf-8", errors="ignore")
    except Exception:
        text = ""
    return [("text", None, text)]


def _parse_pdf(content: bytes) -> List[Tuple[Optional[str], Optional[int], str]]:
    try:
        from pypdf import PdfReader
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "pypdf is not installed. Run: pip install -r backend/requirements.txt"
        ) from e

    reader = PdfReader(io.BytesIO(content))
    pieces: List[Tuple[Optional[str], Optional[int], str]] = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        text = text.strip()
        if not text:
            # OCR fallback for scanned pages
            ocr_text = _ocr_pdf_page(content, i)
            if ocr_text:
                pieces.append((f"Page {i + 1} (OCR)", i + 1, ocr_text))
            continue
        pieces.append((f"Page {i + 1}", i + 1, text))
    return pieces


def _ocr_pdf_page(content: bytes, page_index: int) -> str:
    """Best-effort OCR of a single PDF page. Silently returns '' if unavailable."""
    try:
        from pdf2image import convert_from_bytes  # optional
    except Exception:
        return ""
    try:
        import pytesseract
        from PIL import Image  # noqa: F401
    except Exception:
        return ""
    try:
        images = convert_from_bytes(
            content, first_page=page_index + 1, last_page=page_index + 1
        )
        if not images:
            return ""
        return pytesseract.image_to_string(images[0], lang="eng+hin").strip()
    except Exception:
        return ""


def _parse_docx(content: bytes) -> List[Tuple[Optional[str], Optional[int], str]]:
    try:
        import docx
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "python-docx is not installed. Run: pip install -r backend/requirements.txt"
        ) from e
    doc = docx.Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    # Also include table cells so forms/tables are searchable.
    for t_idx, table in enumerate(doc.tables):
        for r_idx, row in enumerate(table.rows):
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                paragraphs.append(f"[Table {t_idx + 1} Row {r_idx + 1}] " + " | ".join(cells))
    text = "\n".join(paragraphs)
    return [("body", None, text)] if text else []


def _parse_pptx(content: bytes) -> List[Tuple[Optional[str], Optional[int], str]]:
    try:
        from pptx import Presentation
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "python-pptx is not installed. Run: pip install -r backend/requirements.txt"
        ) from e
    prs = Presentation(io.BytesIO(content))
    pieces: List[Tuple[Optional[str], Optional[int], str]] = []
    for i, slide in enumerate(prs.slides):
        parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        text = "\n".join(p.strip() for p in parts if p.strip())
        if text:
            pieces.append((f"Slide {i + 1}", i + 1, text))
    return pieces


def _parse_spreadsheet(
    content: bytes, filename: str
) -> List[Tuple[Optional[str], Optional[int], str]]:
    """Turn a spreadsheet into one searchable chunk per row, prefixed with headers."""
    try:
        import pandas as pd
    except Exception as e:  # pragma: no cover
        raise RuntimeError("pandas is required for spreadsheet ingestion") from e

    name = filename.lower()
    pieces: List[Tuple[Optional[str], Optional[int], str]] = []

    def _df_to_pieces(df, sheet_label: Optional[str]):
        if df is None or df.empty:
            return
        headers = [str(c) for c in df.columns]
        # Row summary chunk: aggregate top values so "overview" questions get a hit
        preview_rows = df.head(20).to_dict("records")
        summary_lines = [f"Sheet: {sheet_label or 'default'}", "Columns: " + ", ".join(headers)]
        for r in preview_rows:
            summary_lines.append(
                " | ".join(f"{k}: {v}" for k, v in r.items() if v is not None and str(v) != "nan")
            )
        pieces.append((sheet_label, 0, "\n".join(summary_lines)))

        # Per-row chunks (capped to keep ingest bounded)
        max_rows = min(len(df), MAX_CHUNKS_PER_FILE - len(pieces))
        for idx, row in df.head(max_rows).iterrows():
            row_txt = " | ".join(
                f"{h}: {row[h]}" for h in headers if h in row and row[h] is not None and str(row[h]) != "nan"
            )
            if row_txt:
                pieces.append((sheet_label, int(idx) + 1, row_txt))

    if name.endswith(".csv") or name.endswith(".tsv"):
        sep = "\t" if name.endswith(".tsv") else ","
        try:
            df = pd.read_csv(io.BytesIO(content), sep=sep, on_bad_lines="skip", encoding="utf-8")
        except UnicodeDecodeError:
            df = pd.read_csv(io.BytesIO(content), sep=sep, on_bad_lines="skip", encoding="latin-1")
        _df_to_pieces(df, None)
        return pieces

    # Excel: iterate every sheet
    xls = pd.ExcelFile(io.BytesIO(content))
    for sheet in xls.sheet_names:
        try:
            df = xls.parse(sheet)
        except Exception:
            continue
        _df_to_pieces(df, str(sheet))
    return pieces


def _parse_image(content: bytes) -> List[Tuple[Optional[str], Optional[int], str]]:
    try:
        import pytesseract
        from PIL import Image
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "pytesseract/Pillow not installed. Run: pip install -r backend/requirements.txt "
            "and install Tesseract OCR system binary (with eng+hin language packs)."
        ) from e
    try:
        img = Image.open(io.BytesIO(content))
    except Exception as e:
        raise RuntimeError(f"Could not open image: {e}") from e
    try:
        text = pytesseract.image_to_string(img, lang="eng+hin")
    except pytesseract.TesseractNotFoundError as e:
        raise RuntimeError(
            "Tesseract OCR binary not found. Install from "
            "https://github.com/UB-Mannheim/tesseract/wiki and ensure it is on PATH."
        ) from e
    text = (text or "").strip()
    return [("ocr", None, text)] if text else []


def _parse_sql(content: bytes) -> List[Tuple[Optional[str], Optional[int], str]]:
    try:
        import sqlparse
    except Exception:
        sqlparse = None  # type: ignore
    try:
        text = content.decode("utf-8", errors="ignore")
    except Exception:
        return []
    if not text.strip():
        return []
    if sqlparse is not None:
        statements = [s.strip() for s in sqlparse.split(text) if s.strip()]
    else:
        statements = [s.strip() for s in text.split(";") if s.strip()]
    return [(f"Statement {i + 1}", i + 1, s) for i, s in enumerate(statements)]


_PARSERS = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "txt": _parse_txt,
    "sql": _parse_sql,
    "image": _parse_image,
}


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------


_WS_RE = re.compile(r"[ \t]+")
_NL_RE = re.compile(r"\n{3,}")


def _clean(text: str) -> str:
    text = text.replace("\x00", " ")
    text = _WS_RE.sub(" ", text)
    text = _NL_RE.sub("\n\n", text)
    return text.strip()


def _split_chunks(text: str) -> List[str]:
    """Split long text into ~CHUNK_SIZE char pieces with overlap at sentence boundaries."""
    text = _clean(text)
    if not text:
        return []
    if len(text) <= CHUNK_SIZE:
        return [text]

    chunks: List[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + CHUNK_SIZE, n)
        # Try to break at a newline or sentence end to avoid cutting words.
        if end < n:
            window = text[start:end]
            for sep in ("\n\n", "\n", ". ", "। ", "? ", "! "):
                idx = window.rfind(sep)
                if idx >= CHUNK_SIZE // 2:
                    end = start + idx + len(sep)
                    break
        chunks.append(text[start:end].strip())
        if end >= n:
            break
        start = max(end - CHUNK_OVERLAP, start + 1)
    return [c for c in chunks if c]


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def _file_hash(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


def ingest_document(
    db: Session,
    *,
    case_id: Optional[str],
    file_name: str,
    content: bytes,
    category: Optional[str] = None,
    tags: Optional[List[str]] = None,
    uploaded_by: Optional[str] = None,
    title: Optional[str] = None,
) -> KbDocument:
    """Ingest one document end-to-end. Idempotent on (case_id, file_hash).

    Returns the KbDocument row with status = 'completed' on success or
    'error' with `error_message` populated on failure. The function never
    raises; errors are recorded on the document so the UI can show them.
    """
    fhash = _file_hash(content)

    existing = (
        db.query(KbDocument)
        .filter(KbDocument.case_id == case_id, KbDocument.file_hash == fhash)
        .first()
    )
    if existing:
        return existing

    src = _source_type(file_name)
    doc = KbDocument(
        id=str(uuid.uuid4()),
        case_id=case_id,
        file_name=file_name,
        file_hash=fhash,
        file_size=len(content),
        source_type=src,
        category=category,
        title=title or os.path.splitext(file_name)[0],
        status="processing",
        tags=tags or [],
        chunk_count=0,
        processing_started_at=datetime.utcnow(),
        created_by=uploaded_by,
    )
    db.add(doc)
    db.commit()

    try:
        if src == "unknown":
            raise RuntimeError(f"Unsupported file type: .{_ext(file_name)}")
        if src in ("xlsx", "csv"):
            pieces = _parse_spreadsheet(content, file_name)
        else:
            parser = _PARSERS.get(src)
            if parser is None:
                raise RuntimeError(f"No parser for source type: {src}")
            pieces = parser(content)

        if not pieces:
            raise RuntimeError("No extractable text was found in this document.")

        # Expand pieces into size-bounded chunks while preserving metadata.
        prepared: List[Dict[str, Any]] = []
        for section, page, piece_text in pieces:
            piece_chunks = _split_chunks(piece_text)
            for i, ctext in enumerate(piece_chunks):
                prepared.append(
                    {
                        "section": section,
                        "page": page,
                        "row_index": (page if src in ("xlsx", "csv") else None),
                        "chunk_index": len(prepared),
                        "text": ctext,
                    }
                )
                if len(prepared) >= MAX_CHUNKS_PER_FILE:
                    break
            if len(prepared) >= MAX_CHUNKS_PER_FILE:
                break

        if not prepared:
            raise RuntimeError("Document had content but produced no usable chunks.")

        texts = [p["text"] for p in prepared]
        vectors = emb_svc.embed_texts(texts)

        langs: List[str] = []
        for p, vec in zip(prepared, vectors):
            ents = extract_entities(p["text"])
            lang = detect_language(p["text"])
            langs.append(lang)
            db.add(
                KbChunk(
                    id=str(uuid.uuid4()),
                    document_id=doc.id,
                    case_id=case_id,
                    chunk_index=p["chunk_index"],
                    text=p["text"],
                    page=p["page"],
                    section=p["section"],
                    row_index=p["row_index"],
                    entities=ents,
                    embedding=emb_svc.pack_vector(vec),
                )
            )

        # Aggregate language signal
        if "hi" in langs and "en" in langs or "mixed" in langs:
            doc.language = "mixed"
        elif "hi" in langs:
            doc.language = "hi"
        else:
            doc.language = "en"
        doc.chunk_count = len(prepared)
        doc.status = "completed"
        doc.processing_completed_at = datetime.utcnow()
        db.commit()
    except Exception as e:
        db.rollback()
        # Reload document in new transaction to mark as error
        row = db.query(KbDocument).filter(KbDocument.id == doc.id).first()
        if row is not None:
            row.status = "error"
            row.error_message = str(e)[:500]
            row.processing_completed_at = datetime.utcnow()
            db.commit()
            doc = row
    return doc
